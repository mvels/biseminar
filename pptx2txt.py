from tutorial.DataModel import Lecture, db
import peewee
import os.path
from pptx import Presentation, exc

class Pptx2Txt(object):
    def __init__(self):
        pass

    def convert(self, ifile, ofile=None):
        retval = ''
        try:
            prs = Presentation(ifile)

            # text_runs will be populated with a list of strings,
            # one for each text run in presentation
            text_runs = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)

            retval = ' '.join(text_runs)

        except exc.PythonPptxError as e:
            print "Could not extract text {0}".format(e)
        return retval


def main():
    pptx2txt = Pptx2Txt()
    lectures = Lecture.select().where(Lecture.content == '', Lecture.url % "*pptx").limit(1000)
    for lecture in lectures:
        if not os.path.exists(lecture.path):
            print "File not found: {0}".format(lecture.path)
            continue
        print lecture.url
        lecture.content = pptx2txt.convert(lecture.path)
        try:
            with db.transaction():
                lecture.save()
        except peewee.OperationalError as e:
            print e

if __name__ == '__main__':
    main()
